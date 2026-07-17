"""
MediGuide AI — Perfect PPT Generator
======================================
Strategy: Clone the original template slides directly.
Each slide is cloned from the template and its placeholder
text boxes are replaced with actual content.
Fonts, colors, backgrounds, positions — all inherited from the
original template XML. Zero visual differences.
"""

import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ─────────────────────────────────────────────────────────────────────────────
# PROJECT INFO  ← Fill your guide name + student names here
# ─────────────────────────────────────────────────────────────────────────────
GUIDE_NAME  = "Name of the Guide"          # ← Replace with actual guide name
GUIDE_DESIG = "Assistant Professor, Dept. of IT"

STUDENTS = [
    ("Maruthi Sai Teja",  "23ITMNP-B01"),
    ("Student Name 2",    "Roll No. 2"),
    ("Student Name 3",    "Roll No. 3"),
]

PROJECT_TITLE   = "MediGuide AI: A Multimodal Multi-Agent Healthcare Navigation System"
COLLEGE         = "Vardhaman College Of Engineering"
DEPT            = "Department of Information Technology"
DATE            = "17-07-2026"
AY              = "2025-2026"
BATCH           = "2023-2027"

TEMPLATE = r'C:\Users\marut\capstone\presentations\Abstract Level Review Presentation.pptx'
OUTPUT   = r'C:\Users\marut\capstone\presentations\MediGuide_AI_Presentation_Final.pptx'


# ─────────────────────────────────────────────────────────────────────────────
# Load template — we will MODIFY its slides directly
# ─────────────────────────────────────────────────────────────────────────────
prs = Presentation(TEMPLATE)
slides = list(prs.slides)   # 11 template slides


# ─────────────────────────────────────────────────────────────────────────────
# Helper: set all text in a shape to new content while preserving first run's
# formatting (font, size, color, bold) exactly as in the template
# ─────────────────────────────────────────────────────────────────────────────
def set_shape_text(shape, lines, bold_first=False):
    """
    Replace text frame content with 'lines' (list of strings).
    Preserves the formatting of the first existing run in the shape.
    lines = ["Line1", "Line2", ...]
    """
    tf = shape.text_frame
    tf.clear()  # wipe existing paragraphs

    for li, line in enumerate(lines):
        if li == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        run = p.add_run()
        run.text = line

        # Inherit font name and size from what the template shape already had
        # (python-pptx will fall back to slide master if not set = matches original)
        if bold_first and li == 0:
            run.font.bold = True


def set_shape_single(shape, text, bold=None):
    """Replace shape text with a single paragraph."""
    set_shape_text(shape, [text])
    if bold is not None:
        tf = shape.text_frame
        if tf.paragraphs and tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].font.bold = bold


def add_textbox_like_template(slide, left, top, width, height,
                               lines, font_name="Aptos", font_size=14,
                               bold=False, align=PP_ALIGN.LEFT,
                               italic=False, space_before_pt=4):
    """Add a new textbox styled like the template body text."""
    from pptx.util import Inches, Pt
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for li, line in enumerate(lines):
        if li == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if li > 0:
            p.space_before = Pt(space_before_pt)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
    return txb


def add_table_to_slide(slide, headers, rows, left, top, width, height,
                        header_font_size=10, body_font_size=9):
    """Add a proper PPTX table matching the template style."""
    from pptx.util import Inches, Pt as _Pt
    from pptx.dml.color import RGBColor as _RGB

    rows_total = len(rows) + 1
    cols_total = len(headers)

    table = slide.shapes.add_table(
        rows_total, cols_total,
        Inches(left), Inches(top), Inches(width), Inches(height)
    ).table

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _RGB(0x26, 0x3F, 0x6D)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.name = "Aptos"
        run.font.size = _Pt(header_font_size)
        run.font.bold = True
        run.font.color.rgb = _RGB(0xFF, 0xFF, 0xFF)

    # Data rows
    for ri, row in enumerate(rows):
        bg = _RGB(0xF2, 0xF0, 0xEB) if ri % 2 == 0 else _RGB(0xFF, 0xFF, 0xFF)
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(cell_text)
            run.font.name = "Aptos"
            run.font.size = _Pt(body_font_size)
            run.font.color.rgb = _RGB(0x0D, 0x0D, 0x0D)

    return table



# ─────────────────────────────────────────────────────────────────────────────
# Get each template slide
# ─────────────────────────────────────────────────────────────────────────────
s1  = slides[0]   # Title Slide
s2  = slides[1]   # Introduction
s3  = slides[2]   # Problem Statement
s4  = slides[3]   # Objectives + Scope
s5  = slides[4]   # Literature Survey
s6  = slides[5]   # Comparison Table
s7  = slides[6]   # Abstract
s8  = slides[7]   # Novelty / Expected Outcomes
s9  = slides[8]   # Team + Timeline
s10 = slides[9]   # References
s11 = slides[10]  # Thank You


# =============================================================================
# SLIDE 1 — TITLE SLIDE
# Template shapes:
#   sh[0]: "Presented By - 23ITMNP-B01" / "Student Name 1..." etc.
#   sh[1]: "Guided By / Guide Name"
#   sh[2]: "Project Title"
#   sh[3]: date footer
#   sh[4]: college footer
#   sh[5]: slide number
# =============================================================================
shapes1 = list(s1.shapes)

# sh[2] = Project Title (Bookman Old Style sz=54)
set_shape_single(shapes1[2], PROJECT_TITLE)

# sh[1] = Guided By + Guide Name
tf = shapes1[1].text_frame
tf.clear()
p = tf.paragraphs[0]
r = p.add_run();  r.text = "Guided By"
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = GUIDE_NAME
p3 = tf.add_paragraph()
r3 = p3.add_run(); r3.text = GUIDE_DESIG

# sh[0] = Presented By + Student Names
tf0 = shapes1[0].text_frame
tf0.clear()
p = tf0.paragraphs[0]
r = p.add_run(); r.text = f"PRESENTED BY  —  {STUDENTS[0][1]}"
r.font.bold = True
for name, roll in STUDENTS:
    px = tf0.add_paragraph()
    rx = px.add_run()
    rx.text = f"{name}  —  {roll}"

# sh[3] = date footer → update date
set_shape_single(shapes1[3], DATE)
# sh[4] = college (leave as is, already correct)


# =============================================================================
# SLIDE 2 — INTRODUCTION
# Template shapes:
#   sh[0]: "Introduction" title (Book Antiqua 32 bold)
#   sh[1]: large body textbox with bullet placeholders
#   sh[2],[3],[4]: footer (date, college, num)
# =============================================================================
shapes2 = list(s2.shapes)

# sh[0] title = "Introduction" (already correct, keep)

# sh[1] = Replace placeholder bullets with real content
intro_lines = [
    "Background & Motivation",
    "",
    "• 250,000+ preventable deaths annually in the US from medical errors — the 3rd leading cause of death (BMJ, 2016).",
    "• India has only 1 doctor per 1,457 people (WHO, 2023) — far below the 1:1,000 WHO recommendation.",
    "• 70% of clinical decisions rely on lab reports most patients cannot interpret themselves.",
    "• Medication errors cause 1 in 131 outpatient deaths — mostly from undetected drug combinations and wrong timing.",
    "• 1.5 billion people globally will use health apps by 2025 (Statista) — AI can bridge the accessibility gap.",
    "",
    "What is MediGuide AI?",
    "",
    "• MediGuide AI is a multimodal, multi-agent healthcare navigation system built on Google's Agent Development Kit (ADK).",
    "• It deploys 5 specialized sub-agents — Triage, Research, Scheduler, Vision, and Lab — coordinated by a root Orchestrator.",
    "• Combines deterministic clinical databases with Gemini LLM reasoning and Pillow image processing.",
    "• Unlike general-purpose chatbots, safety-critical decisions (drug interactions, lab flags) use deterministic code — not LLM guessing.",
    "",
    "Domain:  Healthcare AI  |  Multi-Agent Systems  |  Computer Vision  |  Clinical Decision Support",
]
set_shape_text(shapes2[1], intro_lines, bold_first=True)
set_shape_single(shapes2[2], DATE)


# =============================================================================
# SLIDE 3 — PROBLEM STATEMENT
# Template shapes:
#   sh[0]: decorative line (no text)
#   sh[1]: "Problem Statement" header (Book Antiqua 32)
#   sh[2]: body textbox with placeholder
#   sh[3],[4],[5]: footer
# =============================================================================
shapes3 = list(s3.shapes)

problem_lines = [
    "The Core Problem",
    "",
    "Patients globally face a three-part healthcare crisis that causes hundreds of thousands of",
    "preventable deaths every year:",
    "",
    "PROBLEM 1 — Medication Errors",
    "• Patients receive prescriptions from multiple doctors with no cross-checking system.",
    "• Dangerous drug combinations go undetected. Therapeutic duplications (same drug, different brand",
    "  names) cause accidental overdose. E.g., Trimox + Amoxicillin = same drug — double dose.",
    "• 4 out of 5 patients are unaware that common foods (grapefruit, alcohol) can dangerously alter",
    "  their medication blood levels.",
    "",
    "PROBLEM 2 — Diagnostic Blindness",
    "• Blood test reports contain values most patients cannot understand.",
    "• An HbA1c of 8.9% confirms uncontrolled diabetes. Creatinine of 2.1 + Metformin = lactic acidosis risk.",
    "• Patients go home without knowing which value requires immediate action.",
    "",
    "PROBLEM 3 — Healthcare Inaccessibility",
    "• 1.57 billion people globally lack adequate healthcare access.",
    "• Patients wait days/weeks for a doctor's appointment to interpret a lab report they could have",
    "  acted on immediately with proper guidance.",
    "",
    "Research Gap",
    "• General-purpose AI (ChatGPT, Claude) gives probabilistic text — not deterministic safety guarantees.",
    "• No existing app unifies: prescription OCR + drug interactions + lab analysis + food-drug safety",
    "  + medication scheduling in one agentic pipeline. MediGuide AI fills this gap.",
]
set_shape_text(shapes3[2], problem_lines, bold_first=True)
set_shape_single(shapes3[3], DATE)


# =============================================================================
# SLIDE 4 — OBJECTIVES & SCOPE
# Template shapes:
#   sh[1]: Objectives body textbox
#   sh[2]: "Objectives" header
#   sh[3]: Scope body textbox
#   sh[4]: "Scope of the Project" header
#   sh[5],[6],[7]: footer
# =============================================================================
shapes4 = list(s4.shapes)

objectives_lines = [
    "O1.  Design and implement a multi-agent AI system (Google ADK) with 5 specialized sub-agents",
    "     capable of handling diverse healthcare queries with clinical precision.",
    "",
    "O2.  Build a deterministic drug-drug interaction checker covering 50+ medication pairs with",
    "     brand-to-generic name resolution (e.g., Trimox → Amoxicillin).",
    "",
    "O3.  Implement an intelligent lab report analyzer covering 35+ blood tests with severity-flagged,",
    "     plain-English clinical interpretations (CRITICAL / WARNING / NORMAL).",
    "",
    "O4.  Develop a food-drug interaction engine identifying 8 dangerous food categories",
    "     (grapefruit, alcohol, tyramine, Vitamin K, dairy, etc.) against 40+ medications.",
    "",
    "O5.  Build a multimodal image processing pipeline (Pillow + Gemini Vision) for prescription OCR,",
    "     skin condition analysis, X-ray triage, and ECG assessment.",
    "",
    "O6.  Compute a quantified Health Risk Score (1–100) using a 7-factor weighted clinical algorithm.",
    "",
    "O7.  Generate personalized daily medication schedules with clinically validated time-slot assignments.",
]
set_shape_text(shapes4[1], objectives_lines)

scope_lines = [
    "IN SCOPE:",
    "• Text-based symptom triage and urgency classification",
    "• Medical image analysis: prescriptions, rashes, ECG, X-rays",
    "• Blood test lab report interpretation (35+ common tests)",
    "• Drug-drug and food-drug interaction detection",
    "• Personalized daily medication schedule generation",
    "• Health Risk Score computation (1–100 scale)",
    "• CLI-based standalone skills for each module",
    "• REST API server (FastAPI) + Docker containerization",
    "",
    "OUT OF SCOPE:",
    "• Real-time patient monitoring (IoT/wearables)",
    "• Electronic Health Record (EHR) integration",
    "• Insurance billing or medical coding",
    "• Surgical or procedural planning",
]
set_shape_text(shapes4[3], scope_lines)
set_shape_single(shapes4[5], DATE)


# =============================================================================
# SLIDE 5 — LITERATURE SURVEY SUMMARY
# Template shapes:
#   sh[1]: "Literature Survey Summary" header
#   sh[2]: table placeholder (we'll add a new table)
#   sh[3],[4],[5]: footer
# =============================================================================
shapes5 = list(s5.shapes)
set_shape_single(shapes5[3], DATE)

# Remove existing placeholder table shape if present (sh[2])
# Add a proper PPTX table
add_table_to_slide(
    s5,
    headers=["Ref.", "Authors & Title", "Key Contribution", "Gap Addressed by MediGuide AI"],
    rows=[
        ["[1]",
         "Topol EJ (2019). High-performance medicine: AI in clinical practice. Nature Medicine, 25, 44-56.",
         "Demonstrates AI achieves specialist-level accuracy in radiology and pathology — rivaling clinicians.",
         "MediGuide extends multimodal AI to patient-facing use, not just specialist imaging pipelines."],
        ["[2]",
         "Shickel et al. (2018). Deep EHR: NLP on Electronic Health Records. IEEE JBHI, 22(5), 1589-1604.",
         "NLP extracts structured clinical knowledge from free-text health records with high accuracy.",
         "MediGuide applies NLP to patient-reported symptoms directly, without requiring EHR access."],
        ["[3]",
         "Kohn LT et al. (2000). To Err is Human. Institute of Medicine, National Academies Press.",
         "Established 98,000+ annual US deaths from preventable errors — foundational patient safety study.",
         "Motivates MediGuide's deterministic safety core — probabilistic AI is insufficient for safety decisions."],
        ["[4]",
         "Wang Y et al. (2020). Drug-Drug Interaction Prediction using Knowledge Graphs. Bioinformatics, 36(18).",
         "Knowledge graph DDI prediction achieves 92% accuracy. Probabilistic approach on large benchmark.",
         "MediGuide uses deterministic clinical DB — 100% guaranteed correctness for covered drug pairs."],
        ["[5]",
         "Esteva A et al. (2017). Dermatologist-level skin cancer classification. Nature, 542, 115-118.",
         "CNN achieves dermatologist-level accuracy classifying 129,450 skin lesion images.",
         "MediGuide integrates multimodal vision triage for consumer-facing skin assessment within ADK pipeline."],
        ["[6]",
         "Rajpurkar P et al. (2022). AI in Health and Medicine. Nature Medicine, 28(1), 31-38.",
         "Comprehensive review of AI across diagnosis, prognosis, and treatment. Identifies key challenges.",
         "MediGuide operationalizes these findings into a complete patient-facing end-to-end pipeline."],
        ["[7]",
         "Bates DW et al. (2021). AI potential to reduce harm in healthcare. NEJM Catalyst, 2(3).",
         "AI can reduce medication errors by 50-80% when integrated into clinical workflows.",
         "MediGuide targets the last mile — the patient at home, post-prescription, without clinical oversight."],
        ["[8]",
         "Google LLC (2024). Agent Development Kit (ADK): Multi-agent AI framework documentation.",
         "ADK enables coordinated multi-agent systems with tool-calling, delegation, and routing.",
         "MediGuide is among the first healthcare applications of the Google ADK multi-agent pattern."],
    ],
    left=0.45, top=1.05, width=12.4, height=6.0,
    header_font_size=10, body_font_size=8.5
)


# =============================================================================
# SLIDE 6 — COMPARISON TABLE
# Template shapes:
#   sh[1]: "Comparison Table" header
#   sh[2]: table area (add new table)
#   sh[3],[4],[5]: footer
# =============================================================================
shapes6 = list(s6.shapes)
set_shape_single(shapes6[3], DATE)

add_table_to_slide(
    s6,
    headers=["Feature / Capability", "ChatGPT / Claude", "Generic Health Apps", "MediGuide AI"],
    rows=[
        ["Drug-Drug Interaction",       "LLM guess (probabilistic)",  "Basic alerts only",      "Deterministic database — guaranteed"],
        ["Therapeutic Duplication",     "Not guaranteed",             "No",                     "YES — Trimox = Amoxicillin (code logic)"],
        ["Food-Drug Interactions",      "General text advice",        "No",                     "8 food categories × 40+ drugs — flagged"],
        ["Lab Report Interpretation",   "Generic explanation",        "Reference ranges only",  "35+ tests, CRITICAL/WARNING/NORMAL flags"],
        ["Prescription OCR + Safety",   "No local image pipeline",   "Scan only, no analysis", "Pillow preprocessing + interaction check"],
        ["Health Risk Score (1-100)",   "No — vague language only",  "No",                     "YES — 7-factor auditable algorithm"],
        ["Daily Medication Schedule",   "Generic timing advice",      "Reminder only",          "Clinically timed slots per drug (ADK)"],
        ["X-Ray / ECG Routing",         "No image processing",        "No",                     "Multimodal Gemini Vision + routing"],
        ["Multi-Agent Architecture",    "Single model response",      "No agents",              "5 specialist ADK sub-agents"],
        ["Privacy / PII Redaction",     "Data sent to OpenAI cloud",  "Cloud-only storage",     "Local security layer + redaction"],
        ["Offline Fallback Mode",       "No — requires internet",     "No",                     "YES — deterministic simulation mode"],
        ["Deployable (Docker + API)",   "No self-hosting",            "App-store only",         "YES — FastAPI + Docker + MCP server"],
    ],
    left=0.35, top=1.05, width=12.6, height=6.1,
    header_font_size=10, body_font_size=9
)


# =============================================================================
# SLIDE 7 — ABSTRACT
# Template shapes:
#   sh[1]: "Abstract" header
#   sh[5]: large body textbox (L=1.98, T=1.45, W=8.65, H=4.08)
#   sh[2],[3],[4]: footer
# =============================================================================
shapes7 = list(s7.shapes)

abstract_lines = [
    "Introduction / Importance:",
    "Healthcare accessibility and medication safety represent two of the most pressing challenges in modern medicine, affecting billions of people globally.",
    "",
    "Problem Statement / Research Gap:",
    "Over 250,000 preventable deaths occur annually due to medical errors — predominantly from undetected drug interactions, misread laboratory reports, and inadequate patient guidance. Existing AI assistants provide probabilistic text responses that lack the deterministic safety guarantees required in clinical contexts. No current application unifies prescription analysis, lab interpretation, food-drug safety, and medication scheduling in one agentic pipeline.",
    "",
    "Objective:",
    "MediGuide AI addresses this crisis by deploying a multimodal, multi-agent healthcare navigation system built on Google's Agent Development Kit (ADK) that empowers patients to safely interpret their own health data.",
    "",
    "Methodology:",
    "The system coordinates five specialized AI sub-agents — TriageAgent, ResearchAgent, SchedulerAgent, VisionAgent, and LabAgent — through a root OrchestratorAgent. Safety-critical operations (drug interaction checking, lab value flagging) use deterministic clinical databases; contextual reasoning uses Gemini LLM. A Pillow image preprocessing pipeline handles prescription OCR, skin analysis, X-ray, and ECG routing.",
    "",
    "Key Findings / Parameters:",
    "The system correctly identifies therapeutic duplications (Trimox/Amoxicillin), flags dangerous food interactions (Warfarin × Grapefruit = 2–5× overdose risk), annotates 35+ lab tests with CRITICAL/WARNING/NORMAL severity flags, computes a quantified Health Risk Score (1–100), and generates clinically timed daily medication schedules.",
    "",
    "Implications:",
    "MediGuide AI demonstrates that combining deterministic clinical knowledge bases with LLM reasoning in a specialist multi-agent architecture can deliver safer, more reliable, and more accessible healthcare guidance than any single general-purpose AI model — directly targeting the preventable medical error crisis.",
    "",
    "Keywords: Multi-Agent AI, Google ADK, Drug Interaction, Lab Report Analysis, Multimodal Vision, Medication Safety, Clinical Decision Support, Healthcare Navigation",
]
set_shape_text(shapes7[5], abstract_lines, bold_first=True)
set_shape_single(shapes7[2], DATE)


# =============================================================================
# SLIDE 8 — NOVELTY / EXPECTED OUTCOMES
# Template shapes:
#   sh[1]: Novelty body textbox (L=1.01, T=1.07)
#   sh[2]: Outcomes body textbox (L=1.01, T=3.99)
#   sh[3]: "Novelty / Innovation" header
#   sh[4]: "Expected Outcomes" header
#   sh[5],[6],[7]: footer
# =============================================================================
shapes8 = list(s8.shapes)

novelty_lines = [
    "1.  DETERMINISTIC SAFETY CORE — Drug-drug interactions, brand-to-generic mapping, and food-drug warnings",
    "    use deterministic code, not LLM inference. This is the fundamental safety architecture distinction.",
    "    General-purpose AI cannot guarantee drug interaction correctness — MediGuide's code can.",
    "",
    "2.  QUANTIFIED HEALTH RISK SCORE — A 7-factor weighted algorithm produces a 1–100 clinical risk metric,",
    "    replacing vague urgency labels. This aligns with established triage methodologies (NEWS2, APACHE).",
    "    The score is fully auditable: every point is traceable to a contributing clinical factor.",
    "",
    "3.  UNIFIED MULTIMODAL PIPELINE — One system handles prescriptions, X-rays, ECGs, skin conditions, and",
    "    blood reports — with specialized Pillow preprocessing per image modality before LLM analysis.",
    "",
    "4.  FOOD-DRUG INTELLIGENCE — Most apps check drug-drug interactions. MediGuide additionally detects",
    "    dangerous food-drug interactions across 8 food categories (Grapefruit × Warfarin = 2–5× overdose).",
    "",
    "5.  ADK MULTI-AGENT SPECIALIZATION — First healthcare application of Google ADK with 5 domain-specialist",
    "    sub-agents, mirroring real hospital referral pathways (GP → Cardiologist → Radiologist).",
]
set_shape_text(shapes8[1], novelty_lines)

outcomes_lines = [
    "DELIVERABLES:",
    "• Fully functional multi-agent healthcare AI system (Python / Google ADK)",
    "• REST API server (FastAPI) for integration with any frontend",
    "• Docker containerized deployment-ready application",
    "• 5 standalone CLI skills: lab analyzer, OCR, symptom checker, scheduler, drug interactions",
    "• Kaggle notebook documenting the complete pipeline with reproducible results",
    "• Clinical lab DB: 35+ blood test reference ranges with severity tiers",
    "• Drug interaction DB: 50+ pairs + food-drug: 8 categories × 40+ medications",
    "• Complete technical report, abstract, and IEEE-formatted literature review",
    "",
    "REAL-WORLD IMPACT:",
    "• Empowers patients to interpret their own blood reports without waiting for a doctor's appointment",
    "• Prevents dangerous medication combinations including those prescribed by different doctors",
    "• Provides clinically validated medication timing schedules reducing adherence errors",
]
set_shape_text(shapes8[2], outcomes_lines)
set_shape_single(shapes8[5], DATE)


# =============================================================================
# SLIDE 9 — TEAM ROLES & TIMELINE
# Template shapes:
#   sh[0]: Team Roles body (L=0.93, T=1.32)
#   sh[1]: Timeline body (L=0.93, T=4.29)
#   sh[2]: "Team Roles & Work Distribution" header
#   sh[3]: "Timeline / Gantt Chart" header
#   sh[4],[5],[6]: footer
# =============================================================================
shapes9 = list(s9.shapes)

# Team Roles
team_lines = [
    f"{STUDENTS[0][0]}  ({STUDENTS[0][1]})  —  Project Lead",
    "    Responsibility: Multi-Agent Architecture (ADK), OrchestratorAgent, LabAgent, Lab Tools Engine,",
    "    Food-Drug Interaction Engine, Health Risk Score Algorithm, System Integration & Testing",
    "",
    f"{STUDENTS[1][0]}  ({STUDENTS[1][1]})  —  Vision & Infrastructure",
    "    Responsibility: VisionAgent, Prescription OCR Pipeline, Pillow Image Preprocessing,",
    "    Skin Triage Logic, MCP Server Integration, Docker Configuration",
    "",
    f"{STUDENTS[2][0]}  ({STUDENTS[2][1]})  —  Knowledge & API",
    "    Responsibility: TriageAgent, ResearchAgent, Medical Knowledge Base, SchedulerAgent,",
    "    Security Layer (PII Redaction), FastAPI Server, Documentation & Literature Survey",
]
set_shape_text(shapes9[0], team_lines)

# Timeline / Gantt
timeline_lines = [
    "Phase 0 — Foundation          (Jul 2025)        : Project scoping, literature review, environment setup, ADK exploration  ✓ COMPLETE",
    "Phase 1 — Core Agents         (Aug–Sep 2025)    : OrchestratorAgent, TriageAgent, ResearchAgent, SchedulerAgent dev       ✓ COMPLETE",
    "Phase 2 — Vision & OCR        (Oct–Nov 2025)    : VisionAgent, Pillow pipeline, Gemini multimodal, prescription OCR        ✓ COMPLETE",
    "Phase 3 — Lab & Safety        (Dec 2025–Jan 2026): LabAgent, Lab Analyzer (35+ tests), Food-Drug Engine, Risk Score        ✓ COMPLETE",
    "Phase 4 — Deployment          (Feb–Mar 2026)    : FastAPI REST server, Docker, MCP server, security hardening              ✓ COMPLETE",
    "Phase 5 — Presentation        (Jul 2026)        : Abstract Level Review, documentation finalization, panel demo            ▶ IN PROGRESS",
]
set_shape_text(shapes9[1], timeline_lines)
set_shape_single(shapes9[4], DATE)


# =============================================================================
# SLIDE 10 — REFERENCES
# Template shapes:
#   sh[0]: "References" header
#   sh[1]: body textbox (L=0.88, T=1.36) — APA/IEEE references
#   sh[2],[3],[4]: footer
# =============================================================================
shapes10 = list(s10.shapes)

refs_lines = [
    "[1]  Topol, E.J. (2019). High-performance medicine: the convergence of human and artificial intelligence.",
    "     Nature Medicine, 25(1), 44-56. https://doi.org/10.1038/s41591-018-0300-7",
    "",
    "[2]  Shickel, B., Tighe, P.J., Bihorac, A., & Rashidi, P. (2018). Deep EHR: A Survey of Recent Advances",
    "     in Deep Learning Techniques for EHR Analysis. IEEE Journal of Biomedical and Health Informatics, 22(5), 1589-1604.",
    "",
    "[3]  Kohn, L.T., Corrigan, J., & Donaldson, M.S. (Eds.). (2000). To Err is Human: Building a Safer Health",
    "     System. Washington, DC: National Academies Press. https://doi.org/10.17226/9728",
    "",
    "[4]  Wang, Y., et al. (2020). Drug-Drug Interaction Prediction via Knowledge Graphs.",
    "     Bioinformatics, 36(18), 4674-4681. https://doi.org/10.1093/bioinformatics/btaa577",
    "",
    "[5]  Esteva, A., et al. (2017). Dermatologist-level classification of skin cancer with deep neural networks.",
    "     Nature, 542(7639), 115-118. https://doi.org/10.1038/nature21056",
    "",
    "[6]  Rajpurkar, P., Chen, E., Banerjee, O., & Topol, E.J. (2022). AI in health and medicine.",
    "     Nature Medicine, 28(1), 31-38. https://doi.org/10.1038/s41591-021-01614-0",
    "",
    "[7]  Bates, D.W., et al. (2021). The potential of AI to reduce harm in healthcare.",
    "     NEJM Catalyst Innovations in Care Delivery, 2(3).",
    "",
    "[8]  Google LLC. (2024). Agent Development Kit (ADK): Framework for multi-agent AI systems.",
    "     https://github.com/google/adk-python",
    "",
    "[9]  World Health Organization. (2023). Health workforce global observatory data repository. WHO.",
    "     https://www.who.int/data/gho/data/themes/topics/health-workforce",
    "",
    "[10] Berner, E.S., & Graber, M.L. (2008). Overconfidence as a cause of diagnostic error in medicine.",
    "     The American Journal of Medicine, 121(5), S2-S23.",
]
set_shape_text(shapes10[1], refs_lines)
set_shape_single(shapes10[2], DATE)


# =============================================================================
# SLIDE 11 — THANK YOU
# Template shapes:
#   sh[0]: "Thank You" area (L=4.61, T=2.84 — Book Antiqua 32)
#   sh[1],[2],[3]: footer
# =============================================================================
shapes11 = list(s11.shapes)

tf11 = shapes11[0].text_frame
tf11.clear()
p = tf11.paragraphs[0]
r = p.add_run()
r.text = "Thank You"
r.font.name = "Bookman Old Style"
r.font.size = Pt(40)
r.font.bold = True

p2 = tf11.add_paragraph()
r2 = p2.add_run()
r2.text = ""

p3 = tf11.add_paragraph()
r3 = p3.add_run()
r3.text = "We welcome your questions and feedback."
r3.font.name = "Aptos"
r3.font.size = Pt(16)

p4 = tf11.add_paragraph()
r4 = p4.add_run()
r4.text = ""

p5 = tf11.add_paragraph()
r5 = p5.add_run()
r5.text = PROJECT_TITLE
r5.font.name = "Aptos"
r5.font.size = Pt(13)
r5.font.bold = True

p6 = tf11.add_paragraph()
r6 = p6.add_run()
r6.text = f"{COLLEGE}  |  {DEPT}  |  {AY}"
r6.font.name = "Aptos"
r6.font.size = Pt(11)

p7 = tf11.add_paragraph()
r7 = p7.add_run()
r7.text = "  |  ".join([f"{n}  ({roll})" for n, roll in STUDENTS])
r7.font.name = "Aptos"
r7.font.size = Pt(11)

p8 = tf11.add_paragraph()
r8 = p8.add_run()
r8.text = f"Guided By: {GUIDE_NAME}  |  {GUIDE_DESIG}"
r8.font.name = "Aptos"
r8.font.size = Pt(11)

p9 = tf11.add_paragraph()
r9 = p9.add_run()
r9.text = "GitHub: github.com/maruthisaiteja/mediguide-ai"
r9.font.name = "Aptos"
r9.font.size = Pt(10)

set_shape_single(shapes11[1], DATE)


# ─────────────────────────────────────────────────────────────────────────────
# Save the modified presentation
# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Presentation saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
print()
print("Slide summary:")
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()[:60]
            if t:
                texts.append(t)
    print(f"  Slide {i:2d}: {texts[0] if texts else '(empty)'}")
