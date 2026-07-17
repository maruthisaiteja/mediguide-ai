from pptx import Presentation

pptx_path = r'C:\Users\marut\capstone\presentations\Abstract Level Review Presentation.pptx'
prs = Presentation(pptx_path)

print(f'Slide dimensions: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches')
print(f'Total slides in template: {len(prs.slides)}')

for i, slide in enumerate(prs.slides, 1):
    print(f'\n=== Slide {i} ({slide.slide_layout.name}) ===')
    for j, shape in enumerate(slide.shapes):
        pos = ''
        if hasattr(shape, 'left') and shape.left is not None:
            pos = f'L={shape.left.inches:.2f} T={shape.top.inches:.2f} W={shape.width.inches:.2f} H={shape.height.inches:.2f}'
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if txt:
                    fn = sz = bd = 'inh'
                    if para.runs:
                        r = para.runs[0]
                        fn = r.font.name or 'inh'
                        sz = int(r.font.size.pt) if r.font.size else 'inh'
                        bd = r.font.bold
                    print(f'  sh[{j}] {pos}')
                    print(f'    "{txt[:110]}"')
                    print(f'    font={fn} sz={sz} bold={bd}')
